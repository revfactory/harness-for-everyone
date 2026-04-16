#!/usr/bin/env python3
"""하네스 대중화 슬라이드 덱 PPTX 생성 (16:9, 10장)"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BASE = Path("/Users/robin/Downloads/harness-easy")
IMG_DIR = BASE / "_workspace/visual/images"
OUT_PPTX = BASE / "하네스_대중화_슬라이드.pptx"

# Colors (palette)
CREAM = RGBColor(0xFF, 0xF8, 0xF0)
TEXT = RGBColor(0x5C, 0x56, 0x50)
CORAL = RGBColor(0xF4, 0xA3, 0x9C)
SAGE = RGBColor(0xA8, 0xC5, 0xA0)
MUSTARD = RGBColor(0xF2, 0xC9, 0x4C)
LAVENDER = RGBColor(0xC3, 0xB1, 0xE1)
SKY = RGBColor(0x89, 0xCF, 0xF0)
GRAY = RGBColor(0xB8, 0xB0, 0xA8)
SAGE_DARK = RGBColor(0x6B, 0x89, 0x62)

# Slide definition data
SLIDES = [
    {
        "num": 1, "layout": "cover",
        "img": "slide_01.png",
        "title": "AI 팀 운영 매뉴얼, 하네스",
        "subtitle": "일반인도 이해할 수 있는 하네스 엔지니어링",
        "message": None,
        "bullets": ["2026-04-16", "황민호(revfactory@gmail.com)"],
        "notes": "안녕하세요. 오늘 약 10분간 '하네스'라는 개념에 대해 이야기하겠습니다. 하네스라는 말을 처음 들어보신 분도 계실 텐데, 전혀 걱정하지 않으셔도 됩니다. 오늘 이 자리를 마치면 '아, 그게 그거구나'라고 느끼실 수 있도록 준비했습니다.\n\n이 발표의 목표는 단 하나입니다. '하네스가 무엇인지 개념적으로 이해하는 것.' 코딩을 배우거나 기술 실습을 하는 자리가 아닙니다. 레시피를 읽는 것과 요리사 자격증을 따는 것은 전혀 다른 일이잖아요. 오늘은 레시피를 읽는 시간입니다.",
    },
    {
        "num": 2, "layout": "image-text",
        "img": "slide_02.png",
        "title": "AI 하나에게 다 시키면?",
        "message": "혼자 코스 요리를 다 차리면 헤맨다.",
        "bullets": [
            "\"기획서 검토하고 수정하고 보고서 만들어줘\"",
            "결과는? 뒤죽박죽",
            "우리도 혼자서 코스 요리를 다 차리지 않는다",
        ],
        "notes": "혹시 이런 경험 있으신가요? AI 챗봇에게 '이 기획서 검토하고, 수정하고, 보고서도 만들어줘'라고 한 번에 부탁했더니 결과가 뒤죽박죽이었던 적. 이건 AI가 멍청해서가 아닙니다. 우리 사람도 마찬가지예요. 저녁 코스 요리를 혼자서 전채부터 디저트까지 동시에 만들면 어떻게 될까요? 그래서 식당에는 여러 셰프가 있는 겁니다. AI도 똑같습니다.",
    },
    {
        "num": 3, "layout": "image-text",
        "img": "slide_03.png",
        "title": "왜 지금 '하네스 이해'가 필요한가",
        "message": "AI는 팀으로 엮는 시대인데, 우리 개념은 챗봇 한 대에 멈춰 있다.",
        "bullets": [
            "AI는 '한 대 쓰는 시대'에서 '팀으로 엮는 시대'로 변화",
            "만들지 않더라도, 읽을 수는 있어야 한다",
            "이해는 기초 리터러시 — 코딩 몰라도 된다",
        ],
        "notes": "지금 AI 세계에 큰 변화가 일어나고 있습니다. 예전에는 AI를 '한 대' 썼지만 이제 여러 AI를 '팀으로 엮는' 방식으로 바뀌고 있습니다. 문제는 우리 머릿속은 아직 '챗봇 한 대'에 멈춰 있다는 겁니다. 이 격차가 커지면 AI에 대해 과신하거나 막연히 불신하게 됩니다. 만들 필요는 없지만 읽을 수는 있어야 합니다.",
    },
    {
        "num": 4, "layout": "image-text",
        "img": "slide_04.png",
        "title": "호통치는 주방이 아닌, 협업하는 팀 주방",
        "message": "하네스는 우리가 AI 팀에게 역할을 나눠주는 주방 운영 매뉴얼이다.",
        "bullets": [
            "이건 독재 주방이 아니에요",
            "서로 '이 간은 어때?'라고 코멘트하는 팀",
            "이 주방의 운영 매뉴얼 = 하네스",
        ],
        "notes": "핵심 비유입니다. 여기서 말하는 주방은 호통치는 주방이 아닙니다. 서로 '이 간은 어때?'라고 코멘트하며 한 접시를 함께 만들어가는 팀 주방입니다. 이 주방에는 각자 역할을 맡은 셰프들이 있고, 헤드셰프가 조율합니다. 그리고 매일 돌아가게 만드는 '운영 매뉴얼' — 이것이 바로 하네스입니다. 사고파는 앱이 아닙니다. 우리가 직접 짜는 설계 방식이에요.",
    },
    {
        "num": 5, "layout": "image-text",
        "img": "slide_05.png",
        "title": "주방의 역할 = AI의 역할",
        "message": "에이전트, 스킬, 오케스트레이터, 하네스가 각각 어떤 주방 역할인지 안다.",
        "bullets": [
            "에이전트 = AI 셰프 (역할을 맡은 AI)",
            "오케스트레이터 = 헤드셰프 (조율자)",
            "스킬 = 레시피 카드 / 하네스 = 운영 매뉴얼",
        ],
        "notes": "에이전트는 각 파트를 맡은 AI 셰프입니다. 오케스트레이터는 헤드셰프로, 직접 요리하지 않고 주문을 받아 나누고 결과를 모읍니다. 스킬은 레시피 카드입니다 — 구체적 절차서. 그리고 이 모든 것을 묶어둔 운영 매뉴얼이 하네스입니다. 제품이 아니라 설계 방식. 매일 조금씩 고쳐 쓰는 살아있는 매뉴얼입니다.",
    },
    {
        "num": 6, "layout": "image-text",
        "img": "slide_06.png",
        "title": "3분이면 감, 30분이면 그림, 1시간이면 설명",
        "message": "3단계 학습 경로로 누구나 자기 속도에 맞게 이해할 수 있다.",
        "bullets": [
            "입문 (1~10분): 하네스를 한 문장으로 말할 수 있다",
            "중급 (30분): 다섯 요소의 흐름도를 그릴 수 있다",
            "심화 (60~90분): 타인에게 10분간 설명할 수 있다",
        ],
        "notes": "학습은 세 단계입니다. 입문은 1~10분이면 됩니다. 이 시간 안에 하네스가 뭔지 한 문장으로 말할 수 있고, 일상에서 비슷한 상황을 떠올릴 수 있으면 성공입니다. 중급은 30분, 심화는 60~90분입니다. 중요한 건 첫 단계에서 멈춰도 된다는 점입니다. 이해만으로 충분합니다.",
    },
    {
        "num": 7, "layout": "image-text",
        "img": "slide_07.png",
        "title": "이건 아니에요 — 흔한 오해 5가지",
        "message": "가장 흔한 5가지 오해를 미리 바로잡는다.",
        "bullets": [
            "X 'AI 제품이다' → O 설계 방식이다",
            "X 'AI가 스스로 판단' → O 설계된 범위 안에서 선택",
            "X '사람 없이 돌아간다' → O 사람이 짜고 사람이 고친다",
        ],
        "notes": "하네스 오해 5가지를 바로잡습니다. 첫째, 하네스는 제품이 아니라 설계 방식. 둘째, AI는 스스로 판단하지 않고 설계 범위 안에서 확률적으로 선택. 셋째, 사람 없이 돌아가지 않으며 가장 중요한 주어는 '우리'. 넷째, AI는 매번 기억을 잃어 메모리 설계가 필요. 다섯째, 코딩 몰라도 이해 가능.",
    },
    {
        "num": 8, "layout": "image-full",
        "img": "slide_08.png",
        "title": "일상에서 찾는 하네스적 순간",
        "message": "내 일상에서 하네스적 순간을 찾을 수 있다.",
        "bullets": [
            "김장: 절이기·양념·속 — 역할 분담의 팀워크",
            "이사: 포장·운반·배치 — 누군가 전체를 조율",
            "소풍: 장보기·요리·자리 — 체크리스트가 매뉴얼",
        ],
        "notes": "김장, 이사, 소풍. 일상에도 하네스적 구조가 있습니다. 가족이 김장할 때 배추·양념·속 담당으로 나뉘는 것 — 에이전트. 어머니가 '배추 몇 포기는 절이고' 조율하는 것 — 오케스트레이터. 매년 '올해 김장 순서'를 적어둔 메모 — 하네스. 여러분의 일상에서도 비슷한 순간을 하나 찾아보세요.",
    },
    {
        "num": 9, "layout": "image-text",
        "img": "slide_09.png",
        "title": "세 문장으로 기억하세요",
        "message": "하네스의 본질을 세 문장과 자가 점검으로 다진다.",
        "bullets": [
            "하네스는 AI 팀 운영 매뉴얼이다 — 제품이 아니라 설계 방식",
            "사람이 짜고, AI가 실행한다 — 주어는 항상 '우리'",
            "AI 셰프는 매번 기억을 잃는다 — 그래서 메모리 설계가 필요",
        ],
        "notes": "세 문장으로 기억하세요. 1) 하네스는 AI 팀 운영 매뉴얼이다. 2) 사람이 짜고 AI가 실행한다. 3) AI 셰프는 매번 기억을 잃는다. 자가 점검: 하네스를 무엇에 비유했나요? 누가 설계하나요? 내 일상의 하네스적 순간은? 이 세 질문에 답할 수 있으면 오늘의 목표는 달성입니다.",
    },
    {
        "num": 10, "layout": "image-text",
        "img": "slide_10.png",
        "title": "더 알고 싶다면",
        "message": "이해 다음의 길이 열려 있다.",
        "bullets": [
            "이해만으로 충분합니다 — 오늘 여기까지가 입문 완료",
            "더 깊이: 중급(30분), 심화(90분) 경로",
            "만들기는 별도 경로 — 이해와 만들기는 다른 능력",
        ],
        "notes": "마지막 슬라이드입니다. 이해만으로 충분합니다. 오늘 10분을 함께하셨다면 하네스가 뭔지 감을 잡으셨고, 일상에서 비슷한 순간을 떠올리셨을 겁니다. 앞으로 뉴스·회의·대화에서 'AI 에이전트', '하네스' 같은 단어를 '아 그거구나' 잡을 수 있습니다. 더 알고 싶다면 중급(30분), 심화(90분) 경로가 있습니다. 이해와 만들기는 다른 능력입니다. 감사합니다.",
    },
]


def set_text(tf, text, *, size=18, color=TEXT, bold=False, align=None, font="Apple SD Gothic Neo"):
    """Set text with styling"""
    tf.clear()
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if align:
            p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.name = font


def add_bg(slide, width, height):
    """Add cream background rectangle"""
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, width, height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = CREAM
    bg.line.color.rgb = CREAM
    bg.line.width = Pt(0)
    # Send to back
    spTree = bg._element.getparent()
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)


def add_slide_number_badge(slide, num, total, width):
    """Top-right slide number badge"""
    x = width - Inches(1.3)
    y = Inches(0.25)
    w = Inches(1.0)
    h = Inches(0.4)
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    box.fill.solid()
    box.fill.fore_color.rgb = CREAM
    box.line.color.rgb = GRAY
    box.line.width = Pt(1.2)
    tf = box.text_frame
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    set_text(tf, f"{num} / {total}", size=12, color=GRAY, align=PP_ALIGN.CENTER)


def add_footer(slide, num, width, height):
    """Bottom-right footer"""
    box = slide.shapes.add_textbox(width - Inches(4.5), height - Inches(0.4),
                                    Inches(4.2), Inches(0.3))
    tf = box.text_frame
    tf.margin_left = 0
    tf.margin_right = 0
    set_text(tf, f"하네스 대중화 · 슬라이드 {num}/10",
             size=10, color=GRAY, align=PP_ALIGN.RIGHT)


def build_cover_slide(slide, data, width, height):
    """Layout: cover — centered image above title"""
    # Image (centered, top portion)
    img_path = str(IMG_DIR / data["img"])
    img_w = Inches(7.5)
    img_h = Inches(4.2)
    img_x = (width - img_w) / 2
    img_y = Inches(0.5)
    slide.shapes.add_picture(img_path, img_x, img_y, width=img_w, height=img_h)

    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(5.0), width - Inches(2), Inches(1.1))
    tf = title_box.text_frame
    tf.margin_top = 0
    tf.margin_bottom = 0
    set_text(tf, data["title"], size=44, color=TEXT, bold=True, align=PP_ALIGN.CENTER)

    # Subtitle
    if data.get("subtitle"):
        sub_box = slide.shapes.add_textbox(Inches(1), Inches(6.0), width - Inches(2), Inches(0.5))
        tf = sub_box.text_frame
        set_text(tf, data["subtitle"], size=20, color=GRAY, align=PP_ALIGN.CENTER)

    # Meta
    if data.get("bullets"):
        meta_box = slide.shapes.add_textbox(Inches(1), Inches(6.7), width - Inches(2), Inches(0.5))
        tf = meta_box.text_frame
        meta_text = "  ·  ".join(data["bullets"])
        set_text(tf, meta_text, size=14, color=GRAY, align=PP_ALIGN.CENTER)


def build_image_text_slide(slide, data, width, height):
    """Layout: image on left, text on right"""
    # Image (left half)
    img_path = str(IMG_DIR / data["img"])
    img_w = Inches(6.0)
    img_h = Inches(4.5)
    img_x = Inches(0.5)
    img_y = (height - img_h) / 2
    slide.shapes.add_picture(img_path, img_x, img_y, width=img_w, height=img_h)

    # Text area (right half)
    text_x = Inches(7.0)
    text_w = width - text_x - Inches(0.5)

    # Title
    title_box = slide.shapes.add_textbox(text_x, Inches(0.9), text_w, Inches(1.4))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.margin_top = 0
    set_text(tf, data["title"], size=28, color=TEXT, bold=True)

    # Message (sage green highlight)
    if data.get("message"):
        msg_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, text_x, Inches(2.6),
                                          text_w, Inches(0.8))
        msg_box.fill.solid()
        msg_box.fill.fore_color.rgb = RGBColor(0xED, 0xF4, 0xEA)
        msg_box.line.color.rgb = SAGE
        msg_box.line.width = Pt(1.5)
        tf = msg_box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.15)
        tf.margin_right = Inches(0.15)
        tf.margin_top = Inches(0.08)
        tf.margin_bottom = Inches(0.08)
        set_text(tf, data["message"], size=15, color=SAGE_DARK, bold=True)

    # Bullets
    if data.get("bullets"):
        bullet_box = slide.shapes.add_textbox(text_x, Inches(3.7), text_w, Inches(3.5))
        tf = bullet_box.text_frame
        tf.word_wrap = True
        tf.margin_top = 0
        for i, bullet in enumerate(data["bullets"]):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_after = Pt(10)
            run = p.add_run()
            run.text = "✦  " + bullet
            run.font.size = Pt(16)
            run.font.color.rgb = TEXT
            run.font.name = "Apple SD Gothic Neo"


def build_image_full_slide(slide, data, width, height):
    """Layout: center image with text above/below"""
    # Title at top
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), width - Inches(1), Inches(0.9))
    tf = title_box.text_frame
    tf.word_wrap = True
    set_text(tf, data["title"], size=30, color=TEXT, bold=True, align=PP_ALIGN.CENTER)

    # Message
    if data.get("message"):
        msg_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), width - Inches(1), Inches(0.5))
        tf = msg_box.text_frame
        set_text(tf, data["message"], size=16, color=SAGE_DARK, bold=True, align=PP_ALIGN.CENTER)

    # Image (center)
    img_path = str(IMG_DIR / data["img"])
    img_w = Inches(9)
    img_h = Inches(4.0)
    img_x = (width - img_w) / 2
    img_y = Inches(1.9)
    slide.shapes.add_picture(img_path, img_x, img_y, width=img_w, height=img_h)

    # Bullets at bottom (horizontal)
    if data.get("bullets"):
        bullet_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.1), width - Inches(1), Inches(1.0))
        tf = bullet_box.text_frame
        tf.word_wrap = True
        text = "     ✦   ".join(data["bullets"])
        set_text(tf, "✦   " + text, size=13, color=TEXT, align=PP_ALIGN.CENTER)


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    width, height = prs.slide_width, prs.slide_height

    blank_layout = prs.slide_layouts[6]  # blank

    for data in SLIDES:
        slide = prs.slides.add_slide(blank_layout)
        add_bg(slide, width, height)

        layout = data["layout"]
        if layout == "cover":
            build_cover_slide(slide, data, width, height)
        elif layout == "image-full":
            build_image_full_slide(slide, data, width, height)
        else:  # default image-text
            build_image_text_slide(slide, data, width, height)

        add_slide_number_badge(slide, data["num"], 10, width)
        add_footer(slide, data["num"], width, height)

        # Speaker notes
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = data["notes"]

    prs.save(OUT_PPTX)
    print(f"Saved {OUT_PPTX}")


if __name__ == "__main__":
    main()
